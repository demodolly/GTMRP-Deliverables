#!/usr/bin/env python3
"""Apply speaker notes to GTM Performance & Readiness management deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from speaker_notes import SPEAKER_NOTES

DECK = Path("/workspace/GTM Performance & Readiness Project Update - August 2026.pptx")


def apply_speaker_notes(deck_path: Path = DECK) -> int:
    prs = Presentation(str(deck_path))
    applied = 0

    for i, slide in enumerate(prs.slides, 1):
        notes_text = SPEAKER_NOTES.get(i, "")
        if not notes_text:
            continue

        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = notes_text
        p.font.size = Pt(11)
        applied += 1

    prs.save(str(deck_path))
    return applied


def main():
    count = apply_speaker_notes()
    print(f"Applied speaker notes to {count} slides in {DECK}")
    print(f"Total slides: {len(Presentation(str(DECK)).slides)}")


if __name__ == "__main__":
    main()
