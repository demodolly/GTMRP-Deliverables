"""Shared single-slide builder — header, context strip, barriers & stakeholder impacts."""

from __future__ import annotations

from dataclasses import dataclass
from typing import Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
MID = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC0, 0x50, 0x00)
CONTEXT_FILL = RGBColor(0xE8, 0xEF, 0xF4)
RISK_FILL = RGBColor(0xFD, 0xF2, 0xE9)
STAKEHOLDER_FILL = RGBColor(0xFF, 0xF8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCC, 0xCC, 0xCC)


@dataclass
class BarrierStakeholder:
    barrier: str
    stakeholders: str
    impact: str


@dataclass
class SummarySlideContent:
    title: str
    goal: str
    context_heading: str
    context_body: str
    benefits: Sequence[str]
    barrier_intro: str
    barriers: Sequence[BarrierStakeholder]
    management_ask: str
    risk_heading: str = "Barriers & Stakeholder Impacts — Why Delivery Is Constrained"
    benefits_label: str = "Key Business Benefits"


def _set_margin(tf, top=0.04, bottom=0.04, left=0.07, right=0.07):
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)


def _add_run(paragraph, text, *, size=10, bold=False, color=DARK, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return run


def _box(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    return shape


def _barrier_block(tf, item: BarrierStakeholder, *, size=8):
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    p.space_after = Pt(1)
    _add_run(p, "▸ ", size=size, bold=True, color=ACCENT)
    _add_run(p, item.barrier, size=size, bold=True, color=DARK)

    p2 = tf.add_paragraph()
    p2.space_after = Pt(1)
    _add_run(p2, "   Stakeholders: ", size=size - 1, bold=True, color=NAVY)
    _add_run(p2, item.stakeholders, size=size - 1, color=MID)

    p3 = tf.add_paragraph()
    p3.space_after = Pt(4)
    _add_run(p3, "   Impact: ", size=size - 1, bold=True, color=ACCENT)
    _add_run(p3, item.impact, size=size - 1, color=MID)


def build_summary_slide(prs: Presentation, content: SummarySlideContent):
    """Single slide: compact goal/benefits top; barriers & stakeholder impacts dominate."""
    layout_idx = 43 if len(prs.slide_layouts) > 43 else len(prs.slide_layouts) - 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    sw = prs.slide_width
    sh = prs.slide_height
    m = Inches(0.22)

    banner_h = Inches(0.95)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, m, m, sw - 2 * m, banner_h)
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    tf = banner.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_margin(tf, 0.06, 0.05, 0.12, 0.12)

    p0 = tf.paragraphs[0]
    _add_run(p0, content.title, size=18, bold=True, color=WHITE)
    p1 = tf.add_paragraph()
    p1.space_before = Pt(3)
    _add_run(p1, "Goal: ", size=10, bold=True, color=WHITE)
    _add_run(p1, content.goal, size=10, color=WHITE)

    ctx_top = m + banner_h + Inches(0.08)
    ctx_h = Inches(1.55)
    ctx_w = sw - 2 * m

    ctx_box = _box(slide, m, ctx_top, ctx_w, ctx_h, CONTEXT_FILL)
    tfc = ctx_box.text_frame
    tfc.word_wrap = True
    _set_margin(tfc, 0.05, 0.04, 0.1, 0.1)

    h = tfc.paragraphs[0]
    _add_run(h, content.context_heading, size=10, bold=True, color=NAVY)

    p = tfc.add_paragraph()
    p.space_before = Pt(3)
    _add_run(p, content.context_body, size=8, color=MID)

    p = tfc.add_paragraph()
    p.space_before = Pt(5)
    _add_run(p, content.benefits_label + ": ", size=8, bold=True, color=NAVY)
    _add_run(p, "  •  ".join(content.benefits), size=7.5, color=MID)

    risk_top = ctx_top + ctx_h + Inches(0.08)
    risk_h = sh - risk_top - m

    risk_box = _box(slide, m, risk_top, ctx_w, risk_h, RISK_FILL)
    tfr = risk_box.text_frame
    tfr.word_wrap = True
    _set_margin(tfr, 0.06, 0.05, 0.1, 0.1)

    h = tfr.paragraphs[0]
    _add_run(h, "⚠  ", size=13, bold=True, color=ACCENT)
    _add_run(h, content.risk_heading, size=12, bold=True, color=NAVY)

    p = tfr.add_paragraph()
    p.space_before = Pt(2)
    _add_run(p, content.barrier_intro, size=8, italic=True, color=MID)

    barriers = list(content.barriers)
    mid = (len(barriers) + 1) // 2
    left_items = barriers[:mid]
    right_items = barriers[mid:]

    col_gap = Inches(0.15)
    col_w = (ctx_w - col_gap) / 2

    left_col = _box(
        slide,
        m + Inches(0.1),
        risk_top + Inches(0.62),
        col_w - Inches(0.05),
        risk_h - Inches(0.72),
        STAKEHOLDER_FILL,
    )
    tfl = left_col.text_frame
    tfl.word_wrap = True
    _set_margin(tfl, 0.04, 0.03, 0.06, 0.06)
    for item in left_items:
        _barrier_block(tfl, item, size=7.5)

    right_col = _box(
        slide,
        m + Inches(0.1) + col_w + col_gap,
        risk_top + Inches(0.62),
        col_w - Inches(0.05),
        risk_h - Inches(0.72),
        STAKEHOLDER_FILL,
    )
    tfr2 = right_col.text_frame
    tfr2.word_wrap = True
    _set_margin(tfr2, 0.04, 0.03, 0.06, 0.06)
    for item in right_items:
        _barrier_block(tfr2, item, size=7.5)

    p = tfr.add_paragraph()
    p.space_before = Pt(2)
    _add_run(p, "Management ask: ", size=8, bold=True, color=NAVY)
    _add_run(p, content.management_ask, size=7.5, color=MID)

    return slide
