#!/usr/bin/env python3
"""Build the UIF single-slide summary with three-zone layout."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DECK = Path("/workspace/GTM Performance & Readiness Project Update - August 2026.pptx")

# Colours aligned to Cisco-style deck (navy header, light panels)
NAVY = RGBColor(0x1F, 0x4E, 0x79)
DARK = RGBColor(0x33, 0x33, 0x33)
MID = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xC0, 0x50, 0x00)  # warning orange
P1_FILL = RGBColor(0xE8, 0xEF, 0xF4)
P2_FILL = RGBColor(0xE2, 0xEF, 0xE2)
RISK_FILL = RGBColor(0xFD, 0xF2, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCC, 0xCC, 0xCC)


def _set_margin(tf, top=0.05, bottom=0.05, left=0.08, right=0.08):
    tf.margin_top = Inches(top)
    tf.margin_bottom = Inches(bottom)
    tf.margin_left = Inches(left)
    tf.margin_right = Inches(right)


def _add_run(paragraph, text, *, size=10, bold=False, color=DARK, italic=False):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return run


def _add_line(tf, text="", *, size=10, bold=False, color=DARK, space_before=0, space_after=3, level=0):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
    if tf.paragraphs[0].text == "" and text and p == tf.paragraphs[0]:
        pass
    elif not text and p.text == "" and p is tf.paragraphs[0]:
        return p
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if text:
        _add_run(p, text, size=size, bold=bold, color=color)
    return p


def _add_mixed_line(tf, parts, space_after=3):
    """parts: list of (text, bold, size, color)"""
    p = tf.add_paragraph()
    p.space_after = Pt(space_after)
    for text, bold, size, color in parts:
        _add_run(p, text, size=size, bold=bold, color=color)
    return p


def _add_bullet(tf, text, *, size=9, bold_prefix=None, space_after=2):
    p = tf.add_paragraph()
    p.level = 0
    p.space_after = Pt(space_after)
    bullet = "• "
    if bold_prefix:
        _add_run(p, bullet + bold_prefix, size=size, bold=True, color=DARK)
        _add_run(p, text, size=size, bold=False, color=MID)
    else:
        _add_run(p, bullet + text, size=size, color=MID)
    return p


def _rounded_box(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.75)
    return shape


def build_uif_slide(prs: Presentation):
    """Add UIF summary slide using blank layout."""
    layout_idx = 43 if len(prs.slide_layouts) > 43 else len(prs.slide_layouts) - 1
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    sw = prs.slide_width
    sh = prs.slide_height
    m = Inches(0.25)

    # ── Top banner ──────────────────────────────────────────────────────
    banner_h = Inches(1.05)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, m, m, sw - 2 * m, banner_h)
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    tf = banner.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_margin(tf, 0.08, 0.06, 0.15, 0.15)

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    _add_run(p0, "Unified Intelligence Framework (UIF): Marketing Attribution & Intelligence", size=20, bold=True, color=WHITE)

    p1 = tf.add_paragraph()
    p1.space_before = Pt(4)
    _add_run(p1, "Goal: ", size=11, bold=True, color=WHITE)
    _add_run(
        p1,
        "Transitioning from legacy CCID/DTID tagging to industry-standard, automated Workfront "
        "Universal Key attribution without interrupting current business operations.",
        size=11,
        bold=False,
        color=WHITE,
    )

    # ── Content area geometry ───────────────────────────────────────────
    content_top = m + banner_h + Inches(0.12)
    content_h = sh - content_top - m
    gap = Inches(0.12)
    usable_w = sw - 2 * m

    left_w = usable_w * 0.29
    center_w = usable_w * 0.29
    right_w = usable_w * 0.42 - gap

    left_x = m
    center_x = left_x + left_w + gap
    right_x = center_x + center_w + gap

    # ── Phase 1 panel ───────────────────────────────────────────────────
    p1_box = _rounded_box(slide, left_x, content_top, left_w, content_h, P1_FILL)
    tf1 = p1_box.text_frame
    tf1.word_wrap = True
    _set_margin(tf1)

    h = tf1.paragraphs[0]
    _add_run(h, "Phase 1 — Hybrid Foundation", size=12, bold=True, color=NAVY)
    p = tf1.add_paragraph()
    _add_run(p, "(In Progress / Near-Term)", size=9, italic=True, color=MID)
    p = tf1.add_paragraph()
    p.space_before = Pt(6)
    _add_run(p, "Focus: ", size=9, bold=True, color=DARK)
    _add_run(p, "Parallel dual ingestion & standardized data schema.", size=9, color=MID)

    p = tf1.add_paragraph()
    p.space_before = Pt(8)
    _add_run(p, "Key Capabilities:", size=10, bold=True, color=NAVY)
    for item in [
        "Consistent UTM governance across Workfront, Stensul, and URL Builders.",
        "Capture & preservation of utm_id, utm_medium, and utm_source at all lead touchpoints.",
        "Dual ingestion supporting legacy CCID/DTID alongside new UTM parameters.",
    ]:
        _add_bullet(tf1, item, size=8)

    p = tf1.add_paragraph()
    p.space_before = Pt(8)
    _add_run(p, "Business Value:", size=10, bold=True, color=NAVY)
    for prefix, rest in [
        ("Zero Downtime: ", "Seamless transition with no \"big-bang\" cutover or reporting breakage."),
        ("Immediate Attribution: ", "Preserves source UTMs and legacy IDs at initial lead creation."),
        ("Single Source of Truth: ", "Standardized Snowflake publishing schema for reporting teams."),
    ]:
        _add_bullet(tf1, rest, size=8, bold_prefix=prefix)

    # ── Phase 2 panel ───────────────────────────────────────────────────
    p2_box = _rounded_box(slide, center_x, content_top, center_w, content_h, P2_FILL)
    tf2 = p2_box.text_frame
    tf2.word_wrap = True
    _set_margin(tf2)

    h = tf2.paragraphs[0]
    _add_run(h, "Phase 2 — Universal Key Automation", size=12, bold=True, color=NAVY)
    p = tf2.add_paragraph()
    _add_run(p, "(Future State)", size=9, italic=True, color=MID)
    p = tf2.add_paragraph()
    p.space_before = Pt(6)
    _add_run(p, "Focus: ", size=9, bold=True, color=DARK)
    _add_run(p, "End-to-end automation via Workfront IDs across Marketo & Eloqua.", size=9, color=MID)

    p = tf2.add_paragraph()
    p.space_before = Pt(8)
    _add_run(p, "Key Capabilities:", size=10, bold=True, color=NAVY)
    for item in [
        "Workfront auto-generates Channel & Content IDs (eliminates manual generation).",
        "Content IDs embedded across all AEM pages and form submissions.",
        "Identical Universal Key schema across both Marketo and Eloqua.",
    ]:
        _add_bullet(tf2, item, size=8)

    p = tf2.add_paragraph()
    p.space_before = Pt(8)
    _add_run(p, "Business Value:", size=10, bold=True, color=NAVY)
    for prefix, rest in [
        ("Automation: ", "Zero-Latency Deployment Readiness from Workfront through MarTech stack."),
        ("Full-Journey Attribution: ", "Complete content-level insight from touchpoint to form submission."),
        ("Strategic Alignment: ", "Direct path to the Adobe North Star architecture."),
    ]:
        _add_bullet(tf2, rest, size=8, bold_prefix=prefix)

    # ── Right risks panel ───────────────────────────────────────────────
    risk_box = _rounded_box(slide, right_x, content_top, right_w, content_h, RISK_FILL)
    tfr = risk_box.text_frame
    tfr.word_wrap = True
    _set_margin(tfr)

    h = tfr.paragraphs[0]
    _add_run(h, "⚠ ", size=12, bold=True, color=ACCENT)
    _add_run(h, "Key Management Decision & Barrier Summary", size=11, bold=True, color=NAVY)

    p = tfr.add_paragraph()
    p.space_before = Pt(4)
    _add_run(p, "Executive Summary of Risks & Dependencies", size=9, bold=True, italic=True, color=MID)

    p = tfr.add_paragraph()
    p.space_before = Pt(10)
    _add_run(p, "System Dependencies", size=10, bold=True, color=NAVY)

    p = tfr.add_paragraph()
    p.space_after = Pt(2)
    _add_run(p, "• ", size=8, color=MID)
    _add_run(p, "One Cisco CRM Freeze (FY27 H1): ", size=8, bold=True, color=DARK)
    _add_run(p, "SFDC changes paused; seller-facing data delayed until release.", size=8, color=MID)

    _add_bullet(
        tfr,
        "CCID/DTID remains mandatory for SFDC lead creation, MSP, and Last Touch reporting.",
        size=8,
        bold_prefix="Legacy Overlap: ",
    )

    p = tfr.add_paragraph()
    p.space_before = Pt(10)
    _add_run(p, "Operational & Adoption Barriers", size=10, bold=True, color=NAVY)

    last_p = tfr.add_paragraph()
    last_p.space_after = Pt(2)
    _add_run(last_p, "• ", size=8, color=MID)
    _add_run(last_p, "Incomplete Workfront Adoption: ", size=8, bold=True, color=DARK)
    _add_run(
        last_p,
        "Teams outside Workfront still rely on manual CTT legacy processes, capping ROI.",
        size=8,
        color=MID,
    )

    last_p = tfr.add_paragraph()
    last_p.space_after = Pt(2)
    _add_run(last_p, "• ", size=8, color=MID)
    _add_run(last_p, "Architecture & Data Governance Pivots: ", size=8, bold=True, color=DARK)
    _add_run(
        last_p,
        "Pending decisions on Tealium vs. Adobe Capture and unpublished Business-Ready "
        "Datasets (BRD) create rework risk for reporting guidance.",
        size=8,
        color=MID,
    )

    return slide


def replace_uif_slides_in_deck(deck_path: Path = DECK):
    """Replace multi-slide UIF section with the single summary slide."""
    prs = Presentation(str(deck_path))

    # Find UIF section slides to remove (section title + 4 content slides)
    titles_to_remove = []
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else ""
        if title == "Unified Intelligence Framework (UIF)" or title.startswith(
            "Unified Intelligence Framework (UIF) —"
        ):
            titles_to_remove.append(i)

    # Remove in reverse order
    for idx in reversed(titles_to_remove):
        r_id = prs.slides._sldIdLst[idx]  # noqa: SLF001
        prs.part.drop_rel(r_id.rId)
        del prs.slides._sldIdLst[idx]

    # Insert new slide after Programme Overview (slide index 4 = 5th slide, 0-based index 4)
    # After removal, find "Portfolio context" slide index
    insert_after = 4  # default: after overview content slide
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and "Portfolio context" in slide.shapes.title.text:
            insert_after = i
            break

    # Build slide at end then reorder — python-pptx lacks insert; add at end and move
    build_uif_slide(prs)
    new_idx = len(prs.slides) - 1
    sldIdLst = prs.slides._sldIdLst
    el = sldIdLst[new_idx]
    del sldIdLst[new_idx]
    sldIdLst.insert(insert_after + 1, el)

    prs.save(str(deck_path))
    print(f"UIF summary slide inserted at position {insert_after + 2} in {deck_path}")
    print(f"Total slides: {len(prs.slides)}")


def build_standalone(output: Path | None = None):
    base = Presentation(str(DECK))
    trim_count = len(base.slides)
    while len(base.slides) > 0:
        r_id = base.slides._sldIdLst[0]  # noqa: SLF001
        base.part.drop_rel(r_id.rId)
        del base.slides._sldIdLst[0]
    build_uif_slide(base)
    out = output or Path("/workspace/UIF_Summary_Slide.pptx")
    base.save(str(out))
    print(f"Standalone slide saved to {out}")


if __name__ == "__main__":
    build_standalone()
