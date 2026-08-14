#!/usr/bin/env python3
"""Build management update deck from BRD Requirements Register + BRD context."""

from __future__ import annotations

import re
from collections import Counter, defaultdict
from copy import deepcopy
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REGISTER = Path("/workspace/BRD_Requirements_Register.xlsx")
SRC_DECK = Path("/workspace/GTM Performance & Readiness Project Update - August 2026.pptx")
OUT_DECK = Path("/workspace/GTM Performance & Readiness Project Update - August 2026.pptx")

# Short names for section headers
PROJECT_SHORT = {
    "Unified Intelligence Framework (UIF) - Phase 1": "UIF Phase 1",
    "DTID Channel Alignment / Reporting Channel": "DTID Channel Alignment",
    "CTT & CTT Request Portal FY27 Updates": "CTT FY27 Bridge Updates",
    "Workfront & CTT Integration POC": "Workfront–CTT Integration POC",
    "Workfront Business-Ready Dataset (GTMRP)": "Business-Ready Workfront Dataset",
    "Unified Intelligence Framework (UIF) - Phase 2": "UIF Phase 2",
    "CTT Decommission": "CTT Decommission",
}

PROJECT_ORDER = [
    "DTID Channel Alignment / Reporting Channel",
    "Unified Intelligence Framework (UIF) - Phase 1",
    "CTT & CTT Request Portal FY27 Updates",
    "Workfront & CTT Integration POC",
    "Workfront Business-Ready Dataset (GTMRP)",
    "Unified Intelligence Framework (UIF) - Phase 2",
    "CTT Decommission",
]

# Narrative content derived from register + BRDs (expand as register is filled)
PROJECT_NARRATIVES = {
    "DTID Channel Alignment / Reporting Channel": {
        "headline": "FY27 Reporting Channel mapping delivered — foundation for hybrid reporting",
        "roadblocks_extra": [],
        "impact_extra": [
            "Without this mapping, FY27 channel reporting would require complex DTID translation in every report.",
        ],
        "forward": [
            "Maintain DTID audit dependency (full audit as of 09 Jun 2026) through production migration.",
            "Use completed Reporting Channel field as input to CTT FY27 bridge and reporting refresh.",
            "No further build required — focus shifts to adoption monitoring and downstream consumption.",
        ],
    },
    "Unified Intelligence Framework (UIF) - Phase 1": {
        "headline": "Core tagging delivered under hybrid legacy constraints — reporting & publishing still gated",
        "roadblocks_extra": [
            "Adobe North Star transition gated by Reporting Team FY27 refresh schedule — limits visibility of future-state architecture.",
            "Marketo vs Eloqua architecture uncertainty (TEA002/TEA003) — Tealium→Adobe Capture direction requires re-validation.",
            "Not all teams onboarded to Workfront — dual operating model persists.",
            "GTMPR UTM strategy for extended fields not yet defined (UTM002 deferred).",
        ],
        "impact_extra": [
            "Reporting teams cannot fully pivot to UTM/Channel ID until legacy CCID/DTID dependencies are retired.",
            "Incomplete AEM Content ID process (PUB002) delays granular content attribution on new pages.",
            "39 Phase 1 requirements still Not Started — REP/Tray/SFDC bridge work at risk of slipping behind FY27 reporting refresh.",
        ],
        "forward": [
            "Complete PUB002 AEM delivery (Stephen Watts, Sep 4) and publishing process with Christy/Jim.",
            "Resolve TEA002/TEA003 review — decide Tealium vs Adobe Capture path before further Marketo investment.",
            "Prioritize REP/Tray requirements tied to FY27 reporting refresh; sequence with Reporting Team capacity.",
            "Publish GTMPR UTM strategy for extended fields to unblock UTM002 and Workfront capture design.",
            "Continue weekly register updates — target Status on all 55 Phase 1 items by end of Q1 FY27.",
        ],
    },
    "CTT & CTT Request Portal FY27 Updates": {
        "headline": "16 of 26 bridge requirements in progress — keeping CTT viable while Workfront adoption catches up",
        "roadblocks_extra": [
            "CTT must remain live for FY27 — full decommission not planned; bridge scope keeps expanding.",
            "Portal requirements (POR001–POR010) not yet started — request intake still on legacy fields.",
            "Historical Activity/Offer ID audit (CTT006/CTT018) depends on GTMPR audit capacity.",
            "OMS and Tray integration constraints — new fields must not break downstream APIs during bridge.",
        ],
        "impact_extra": [
            "Non-Workfront teams cannot report on FY27 attributes without these CTT/portal changes.",
            "Mixed CCID usage (CTT007) continues to distort channel/initiative alignment in reporting.",
            "Delayed portal alignment means requestors still create IDs without FY27 hierarchy and funnel fields.",
        ],
        "forward": [
            "Finish in-flight Activity ID and Offer ID tool changes (CTT002–CTT014) before portal rollout.",
            "Start POR001–POR010 portal alignment in parallel once Activity ID UI is stable.",
            "Schedule GTMPR historical audit for CTT006/CTT018 with explicit rollback (CTT017).",
            "Document explicit exit criteria for bridge vs North Star target state to prevent permanent dual maintenance.",
        ],
    },
    "Workfront & CTT Integration POC": {
        "headline": "POC not started — critical path to reduce manual ID admin and prove Workfront-first model",
        "roadblocks_extra": [
            "Competing priority on CTT FY27 bridge delivery consumes same dev/integration capacity.",
            "One Cisco CRM (OCC) freeze — no SFDC changes allowed; POC must prove value without CRM impact (BR-04, FR-09).",
            "Adobe North Star architecture not finalized — risk of building interim integration that requires rework (BR-09).",
            "Marketing automation and data lake architecture still unsettled — limits end-to-end validation scope.",
        ],
        "impact_extra": [
            "Workfront teams continue manual Activity/Offer ID administration — adoption friction persists.",
            "No evidence base for scale/pause/stop decision on Workfront–CTT integration (BR-05).",
            "CTT decommission blockers remain undocumented — decommission timeline stays ambiguous.",
            "29 Must Have POC requirements untouched — FY27 measurement continuity proof delayed.",
        ],
        "forward": [
            "Secure dedicated POC sprint capacity separate from CTT FY27 maintenance (even 2–3 week pilot).",
            "Define narrow pilot scope: one activation type, auto ID create + write-back + pilot dataset extract.",
            "Align POC outputs explicitly to North Star field model to minimize rework.",
            "Use POC results at Q2 FY27 decision gate: scale integration, extend bridge, or pause Workfront expansion.",
        ],
    },
    "Workfront Business-Ready Dataset (GTMRP)": {
        "headline": "Foundation dataset not started — blocks seller enrichment and attribution model consumption",
        "roadblocks_extra": [
            "No committed Data Engineering capacity — 19 of 20 Must Have requirements Not Started.",
            "Workfront field definitions still evolving with UIF Phase 1/2 — dataset design keeps shifting.",
            "Reporting Team models depend on stable golden dataset; unclear priority vs other FY27 initiatives.",
            "Adobe North Star and One Cisco CRM roadmaps unsettled — downstream SFDC consumption path unclear.",
        ],
        "impact_extra": [
            "Orchestration and seller teams lack actionable Workfront insights for daily activation decisions.",
            "Attribution models cannot consume standardized channel/content attributes at scale.",
            "Analysts continue ad-hoc joins across Workfront exports — inconsistent metrics and slower reporting.",
            "OPS-001 daily refresh SLA (8 AM ET) not achievable without foundation build.",
        ],
        "forward": [
            "Establish Q1 FY27 foundation sprint: GTMRP001–GTMRP005 + GEN001–GEN003 (golden dataset core).",
            "Lock v1 field list with GTMPR governance — defer Phase 2-only fields to avoid churn.",
            "Pair dataset delivery with UIF Phase 1 REP requirements for consistent attribute definitions.",
            "Assign named Data Engineering owner and weekly checkpoint with Reporting/Orchestration stakeholders.",
        ],
    },
    "Unified Intelligence Framework (UIF) - Phase 2": {
        "headline": "Phase 2 gated — Universal Key and dual-ingestion depend on CRM and reporting readiness",
        "roadblocks_extra": [
            "One Cisco CRM resource constraints — SFDC push intentionally bypassed; DSD002 stores metadata in logs only.",
            "Reporting teams cannot amend models yet — North Star visibility deferred.",
            "Phase 2 requires Phase 1 publishing/AEM Content ID completion (PUB002, MAR001–MAR004).",
            "All 12 requirements Not Started — no staging environment for dual-ingestion validation (DSD003).",
        ],
        "impact_extra": [
            "Cannot achieve 'Zero-Latency Deployment Readiness' — Workfront IDs not aligned across AEM and MAPs.",
            "Lead improvement and content attribution remain limited to Phase 1 legacy-compatible scope.",
            "Tray dual-source logic (DSD001) unbuilt — risk of incorrect enrichment when legacy and Workfront IDs coexist.",
        ],
        "forward": [
            "Do not start Phase 2 build until Phase 1 PUB/TEA review items closed and POC decision made.",
            "Negotiate OCC engagement window for SFDC attribute display (even read-only) to validate seller use cases.",
            "Build DSD003 synchronized staging environment as first Phase 2 milestone — no production CRM changes.",
            "Align Phase 2 kickoff to Reporting Team FY27 refresh completion and published North Star field catalog.",
        ],
    },
    "CTT Decommission": {
        "headline": "Decommission plan documented — execution blocked until Workfront adoption and CRM path are clear",
        "roadblocks_extra": [
            "49 requirements Not Started — decommission touches 15+ integrated systems (Eloqua, SFDC, Tray, Stensul, etc.).",
            "Mandatory Eloqua audit/cleanup phase required before OCID removal (per decommission BRD).",
            "SFDC/OCC changes required for lead creation without CCID — blocked by One Cisco CRM programme.",
            "Competing FY27 priority: CTT bridge updates extend tool life rather than retire it.",
        ],
        "impact_extra": [
            "Legacy CCID/DTID remain mandatory for SFDC lead creation, MSP, and Last Touch reporting.",
            "Continued manual governance across Workfront, URL Builder, Adobe, and Ace Reporting.",
            "Technical debt and dual maintenance cost grow each quarter the bridge extends.",
            "41 Must Have decommission items unscheduled — no integrated cutover plan for FY27.",
        ],
        "forward": [
            "Treat decommission as a programme dependent on POC outcome + OCC roadmap + Eloqua audit completion.",
            "Sequence: (1) Eloqua asset audit, (2) Workfront–CTT POC decision, (3) SFDC change window with OCC, (4) phased system cutover.",
            "Maintain decommission register as living dependency map — update as bridge projects deliver.",
            "Leadership decision needed: confirm FY27 as bridge year vs commit decommission start in FY28.",
        ],
    },
}


def load_register():
    wb = load_workbook(REGISTER, read_only=True)
    ws = wb["Requirements Register"]
    headers = [ws.cell(1, c).value for c in range(1, ws.max_column + 1)]
    by_project = defaultdict(list)
    for r in range(2, ws.max_row + 1):
        row = {h: ws.cell(r, c + 1).value for c, h in enumerate(headers)}
        if row.get("Requirement ID"):
            by_project[row["Project / BRD"]].append(row)
    return by_project


def summarize_project(rows: list[dict]) -> dict:
    status = Counter(str(r.get("Status") or "Not Started") for r in rows)
    total = len(rows)
    tracked = total - status.get("Not Started", 0)
    must_ns = sum(
        1
        for r in rows
        if (r.get("Status") or "Not Started") == "Not Started" and r.get("Priority") == "Must Have"
    )
    barriers = []
    deps = []
    impacts = []
    wins = []
    in_flight = []

    for r in rows:
        rid = str(r.get("Requirement ID") or "")
        st = str(r.get("Status") or "Not Started")
        desc = str(r.get("Requirement Description") or "")[:90]
        b = str(r.get("Barriers / Blockers") or "").strip()
        d = str(r.get("Dependencies") or "").strip()
        imp = str(r.get("Impact on Main Project") or "").strip()
        prog = str(r.get("Program / Initiative") or "")

        if b:
            barriers.append(f"{rid}: {b.split(chr(10))[0][:120]}")
        if d:
            deps.append(f"{rid}: {d[:120]}")
        if imp:
            impacts.append(imp[:200])

        if st == "Complete":
            wins.append(f"{rid} ({prog}): {desc}")
        elif st in ("In Progress", "In Review"):
            in_flight.append(f"{rid} ({prog}): {desc}")

    # dedupe impacts
    unique_impacts = list(dict.fromkeys(impacts))
    return {
        "total": total,
        "tracked": tracked,
        "status": status,
        "must_ns": must_ns,
        "barriers": barriers[:6],
        "deps": deps[:4],
        "impacts": unique_impacts[:4],
        "wins": wins[:6],
        "in_flight": in_flight[:6],
    }


def set_title(slide, title: str, subtitle: str | None = None):
    if slide.shapes.title:
        slide.shapes.title.text = title
    if subtitle is not None:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1 and ph.has_text_frame:
                ph.text = subtitle
                break


def add_bullets(text_frame, items: list[str], level0_size=16, level1_size=14):
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(level0_size)
        p.space_after = Pt(8)


def add_content_slide(prs, layout_idx, title, bullets):
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    set_title(slide, title)
    body = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1 and shape.has_text_frame:
            body = shape.text_frame
            break
    if body is None:
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                body = shape.text_frame
                break
    if body:
        add_bullets(body, bullets)
    return slide


def add_section_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[12]  # Section, Title Only
    slide = prs.slides.add_slide(layout)
    set_title(slide, title, subtitle if subtitle else None)
    return slide


def status_bar(summary: dict) -> str:
    st = summary["status"]
    parts = []
    for label in ["Complete", "In Progress", "In Review", "Deferred", "Cancelled", "Blocked", "Not Started"]:
        if st.get(label):
            parts.append(f"{label}: {st[label]}")
    return " | ".join(parts)


def build_where_we_are(name: str, summary: dict, narrative: dict) -> list[str]:
    short = PROJECT_SHORT.get(name, name)
    bullets = [
        f"{summary['total']} requirements in register — {summary['tracked']} actively tracked, {summary['must_ns']} Must Have still Not Started",
        status_bar(summary),
        narrative["headline"],
    ]
    if summary["wins"]:
        bullets.append("Recent wins:")
        bullets.extend(summary["wins"][:4])
    if summary["in_flight"]:
        bullets.append("In flight:")
        bullets.extend(summary["in_flight"][:4])
    if summary["tracked"] == 0:
        bullets.append("BRD approved — delivery not yet started; requirements captured for planning and dependency mapping.")
    return bullets


def build_roadblocks(name: str, summary: dict, narrative: dict) -> list[str]:
    bullets = []
    if summary["barriers"]:
        bullets.append("From requirements register:")
        bullets.extend(summary["barriers"][:4])
    if summary["deps"]:
        bullets.append("Dependencies:")
        bullets.extend(summary["deps"][:3])
    bullets.append("Programme / prioritization constraints:")
    bullets.extend(narrative["roadblocks_extra"][:5])
    return bullets


def build_impact(name: str, summary: dict, narrative: dict) -> list[str]:
    bullets = []
    if summary["impacts"]:
        bullets.append("Documented cross-project impact:")
        bullets.extend(summary["impacts"][:3])
    bullets.append("If we do not achieve this:")
    bullets.extend(narrative["impact_extra"][:4])
    return bullets


def build_forward(name: str, narrative: dict) -> list[str]:
    return narrative["forward"]


def portfolio_summary(by_project: dict) -> dict:
    total = sum(len(v) for v in by_project.values())
    tracked = sum(
        1
        for rows in by_project.values()
        for r in rows
        if (r.get("Status") or "Not Started") != "Not Started"
    )
    complete = sum(
        1
        for rows in by_project.values()
        for r in rows
        if r.get("Status") == "Complete"
    )
    in_prog = sum(
        1
        for rows in by_project.values()
        for r in rows
        if r.get("Status") in ("In Progress", "In Review")
    )
    return {"total": total, "tracked": tracked, "complete": complete, "in_prog": in_prog}


def trim_slides_after(prs: Presentation, keep: int) -> None:
    """Remove slides beyond index `keep` so re-runs do not duplicate content."""
    while len(prs.slides) > keep:
        r_id = prs.slides._sldIdLst[keep]  # noqa: SLF001 — python-pptx has no public delete API
        prs.part.drop_rel(r_id.rId)
        del prs.slides._sldIdLst[keep]


def main():
    by_project = load_register()
    pf = portfolio_summary(by_project)

    prs = Presentation(str(SRC_DECK))
    # Keep existing intro slides (1–3); append new content
    keep = 3
    trim_slides_after(prs, keep)

    # --- Executive summary ---
    add_section_slide(prs, "Executive Summary", "GTM Performance & Readiness — August 2026")

    add_content_slide(
        prs,
        19,
        "Portfolio snapshot — we are delivering, but pivoting as upstream programmes shift",
        [
            f"{pf['total']} requirements across 7 programmes — {pf['tracked']} tracked in register ({pf['complete']} Complete, {pf['in_prog']} In Progress/In Review)",
            "What is working: DTID Channel Alignment delivered (12/12); UIF Phase 1 core UTM/TEA/Eloqua tagging largely Complete; CTT FY27 bridge actively In Progress (16 items)",
            "What is stuck: 164 requirements still Not Started — concentrated in Business-Ready Dataset, Workfront–CTT POC, Phase 2, and CTT Decommission",
            "Root cause: we are executing a hybrid bridge model while Adobe North Star, One Cisco CRM, and Business-Ready Dataset roadmaps remain unclear or under-resourced",
            "Result: repeated replanning, competing dev capacity, and difficulty closing end-to-end delivery",
            "Ask: explicit prioritization, named owners/capacity for foundation datasets and POC, and published dependency timelines from OCC and Reporting",
        ],
    )

    add_content_slide(
        prs,
        19,
        "Cross-programme dependencies creating roadblocks",
        [
            "One Cisco CRM (OCC): SFDC tracking parameter changes frozen — blocks Phase 2 SFDC push, decommission lead-creation changes, and full North Star CRM integration",
            "Adobe North Star: target architecture and Reporting Team FY27 refresh schedule gate visibility — interim bridge work risks rework if field model shifts",
            "Business-Ready Workfront Dataset (GTMRP): 0/20 started — no golden dataset for attribution models, seller enrichment, or self-service reporting",
            "CTT bridge vs decommission tension: FY27 requires CTT enhancements while decommission BRD assumes retirement — dual investment without decision gate",
            "Workfront adoption gap: not all teams onboarded — legacy CCID/DTID remain mandatory for SFDC, MSP, and Last Touch reporting",
            "Capacity fragmentation: same CTT dev/integration team supports FY27 updates, POC, DTID maintenance, and portal alignment",
        ],
    )

    add_content_slide(
        prs,
        19,
        "Proposed way forward — portfolio level",
        [
            "1. Publish a single FY27 dependency map linking UIF Phase 1 → CTT Bridge → POC → GTMRP → Phase 2 → Decommission with explicit decision gates",
            "2. Ring-fence capacity: minimum viable sprint for GTMRP foundation and Workfront–CTT POC (separate from CTT maintenance)",
            "3. Secure OCC engagement window — even read-only SFDC validation — to unblock Phase 2 and decommission planning",
            "4. Align Reporting Team on North Star field catalog and FY27 refresh dates — stop downstream replanning",
            "5. Continue requirements register as single status source — expand tracking from 44 to all 208 items by programme milestone",
            "6. Leadership decision: confirm FY27 as intentional bridge year with exit criteria vs accelerate decommission start",
        ],
    )

    # --- Per-project sections ---
    for project_name in PROJECT_ORDER:
        rows = by_project.get(project_name, [])
        narrative = PROJECT_NARRATIVES[project_name]
        summary = summarize_project(rows)
        short = PROJECT_SHORT.get(project_name, project_name)

        add_section_slide(prs, short, f"{summary['total']} requirements")

        add_content_slide(prs, 19, f"{short} — Where we are", build_where_we_are(project_name, summary, narrative))
        add_content_slide(prs, 19, f"{short} — Roadblocks", build_roadblocks(project_name, summary, narrative))
        add_content_slide(prs, 19, f"{short} — Impact if not achieved", build_impact(project_name, summary, narrative))
        add_content_slide(prs, 19, f"{short} — Proposal to move forward", build_forward(project_name, narrative))

    # Appendix
    add_section_slide(prs, "Appendix", "Requirements Register")
    add_content_slide(
        prs,
        19,
        "Data source & next steps",
        [
            "Source: BRD_Requirements_Register.xlsx (208 requirements from 7 BRDs) — 44 tracked as of August 2026",
            "Register columns: reference fields (BRD metadata) + tracking fields (Status, Barriers, Dependencies, Impact, Owner, Comments)",
            "This deck auto-refreshes from register data via generate_management_deck.py — re-run as you fill in more rows",
            "Recommended: complete Status and Barriers for Business-Ready Dataset and Workfront–CTT POC first (highest leadership visibility)",
            "Questions for management: capacity allocation, OCC timeline, Reporting/North Star field lock, FY27 bridge vs decommission decision",
        ],
    )

    prs.save(str(OUT_DECK))
    print(f"Saved deck with {len(prs.slides)} slides to {OUT_DECK}")


if __name__ == "__main__":
    main()
